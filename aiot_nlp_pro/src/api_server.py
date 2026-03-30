import os
import sys
from pathlib import Path
from typing import List, Dict, Any, Optional
from uuid import uuid4
from datetime import datetime

# --- 引入你现有的引擎 ---
# 假设 rag_engine.py 和 llm_engine.py 在当前目录或 src 目录下
# 如果目录结构不同，请调整导入路径
try:
    from rag_engine import HybridRagEngine
    from llm_engine import LlmEngine
except ImportError:
    # 适配可能的 src 目录结构
    sys.path.insert(0, str(Path(__file__).parent / "src"))
    from rag_engine import HybridRagEngine
    from llm_engine import LlmEngine

from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
import uvicorn
from fastapi.staticfiles import StaticFiles
from docx import Document as DocxDocument
from pptx import Presentation

# ================= 配置与初始化 =================

app = FastAPI(
    title="AI Teacher RAG API",
    description="基于混合检索和 LLM 的智能问答接口",
    version="1.0.0"
)

GENERATED_ROOT = Path("./generated_outputs")
GENERATED_ROOT.mkdir(parents=True, exist_ok=True)
app.mount("/generated-files", StaticFiles(directory=str(GENERATED_ROOT)), name="generated-files")

# 全局变量存储引擎实例
rag_engine: Optional[HybridRagEngine] = None
llm_engine: Optional[LlmEngine] = None

# 系统配置 (与你 app.py 中的配置保持一致)
SYSTEM_CONFIG = {
    "embedding_model": "BAAI/bge-small-zh-v1.5",
    "model_provider": "ollama",
    "model_name": "qwen2.5:7b-instruct-q4_K_M", 
    "ollama_base_url": "http://localhost:11434",
    "data_directory": "./data",
    "cache_directory": "./data_cache"
}

def initialize_system():
    """初始化 RAG 和 LLM 引擎"""
    global rag_engine, llm_engine
    
    print("🚀 Initializing System...")
    
    # 1. 加载文档
    data_dir = SYSTEM_CONFIG.get("data_directory", "./data")
    raw_docs = []
    data_path = Path(data_dir)
    
    if data_path.exists():
        for file_path in data_path.iterdir():
            if file_path.is_file():
                loader = None
                if file_path.suffix.lower() == '.pdf':
                    loader = PyPDFLoader(str(file_path))
                elif file_path.suffix.lower() == '.docx':
                    loader = Docx2txtLoader(str(file_path))
                elif file_path.suffix.lower() == '.txt':
                    loader = TextLoader(str(file_path), encoding='utf-8')
                
                if loader:
                    try:
                        docs = loader.load()
                        raw_docs.extend(docs)
                    except Exception as e:
                        print(f"⚠️ Failed to load {file_path.name}: {e}")
    
    # 2. 切分文档
    processed_docs = []
    if raw_docs:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=500, 
            chunk_overlap=50,
            separators=["\n\n", "\n", "。", "！", "？", " ", ""]
        )
        processed_docs = splitter.split_documents(raw_docs)
        print(f"✅ Processed {len(processed_docs)} document chunks.")
    
    # 3. 初始化引擎
    rag_engine = HybridRagEngine(SYSTEM_CONFIG)
    if processed_docs:
        rag_engine.ingest_documents(processed_docs, index_path=SYSTEM_CONFIG.get("cache_directory", "./data_cache"))
    else:
        print("⚠️ No documents found to index.")
    
    llm_engine = LlmEngine(SYSTEM_CONFIG)
    print("✅ System Initialization Complete.")

# ================= 数据模型定义 =================

class QueryRequest(BaseModel):
    question: str
    topK: int = 5
    sessionId: Optional[str] = "default_session"
    useCoT: bool = False  # 可选：是否启用思维链

class SourceItem(BaseModel):
    content: str
    source: str
    page: Optional[int] = None

class QueryResponse(BaseModel):
    answer: str
    sources: List[SourceItem]
    sessionId: str
    model_used: str

class CoursewareGenerateRequest(BaseModel):
    sessionId: str
    finalRequirements: str
    outline: Optional[str] = None
    templateId: Optional[int] = None

class CoursewareGenerateResponse(BaseModel):
    taskId: str
    status: str
    progress: int
    outline: str
    pptUrl: str
    docUrl: str
    gameUrl: str
    model_used: str
    llm_mock: bool = False
    ai_structured: bool = False

# ================= API 接口实现 =================

@app.on_event("startup")
async def startup_event():
    """服务启动时自动初始化系统"""
    initialize_system()

@app.post("/api/retrieve_and_answer", response_model=QueryResponse)
async def retrieve_and_answer(request: QueryRequest):
    """
    核心接口：接收问题，检索文档，生成答案
    """
    if not rag_engine or not llm_engine:
        raise HTTPException(status_code=503, detail="System not initialized yet")
    
    if not request.question.strip():
        raise HTTPException(status_code=400, detail="Question cannot be empty")
    
    # 1. 执行检索
    # 这里可以加入动态 Top-K 逻辑，或者直接使用请求中的 topK
    k = request.topK if request.topK > 0 else 5
    context_docs = rag_engine.retrieve(request.question, k=k)
    
    # 2. 构建 Prompt
    system_prompt = """你是一个专业的 AI 教师。请严格依据提供的【上下文】回答问题。
    如果上下文中没有答案，请直接告知用户你不知道，不要编造。
    回答要清晰、准确，并尽量引用原文来源。"""
    
    if request.useCoT:
        system_prompt = """你是一位拥有深厚学科知识的 AI 教师助手。请先进行内部逻辑推导（思维链），然后给出最终答案。
        输出格式要求：
        ### 思考过程
        (简要列出推导步骤)
        ### 最终回答
        (正式答案)
        """
    
    context_text = ""
    if context_docs:
        context_text = "\n\n".join([f"[来源{i+1}]: {doc.page_content}" for i, doc in enumerate(context_docs)])
        full_prompt = f"【上下文】:\n{context_text}\n\n【用户问题】:\n{request.question}"
    else:
        full_prompt = request.question
        system_prompt += "\n注意：未提供相关上下文，请仅依靠你的通用知识回答，并提示用户资料可能不足。"
    
    # 3. 调用 LLM 生成
    try:
        response_text = llm_engine.generate(full_prompt, system_prompt=system_prompt)
        
        # 如果启用了 CoT 但前端不需要显示思考过程，可以在这里做简单的文本清洗
        # (此处保留原始返回，由前端决定是否展示思考过程，或者在此处用正则剥离)
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM Generation Error: {str(e)}")
    
    # 4. 格式化参考来源
    sources_list = []
    for doc in context_docs:
        sources_list.append(SourceItem(
            content=doc.page_content[:200] + "...", # 截取部分内容避免返回过长
            source=os.path.basename(doc.metadata.get('source', 'Unknown')),
            page=doc.metadata.get('page', None)
        ))
    
    return QueryResponse(
        answer=response_text,
        sources=sources_list,
        sessionId=request.sessionId,
        model_used=SYSTEM_CONFIG.get("model_name", "unknown")
    )

def _sanitize_filename(raw: str) -> str:
    safe = "".join(ch if ch.isalnum() or ch in "-_" else "_" for ch in raw.strip())
    return safe[:50] if safe else "courseware"

COURSEWARE_JSON_SYSTEM = """你是资深课程设计师。根据用户需求输出**一个 JSON 对象**（不要 markdown 代码块），字段必须齐全：
{
  "course_title": "课程标题",
  "slides": [
    {"title": "单页标题", "bullets": ["要点1", "要点2", "要点3"], "speaker_notes": "讲师备注，可为空字符串"}
  ],
  "lesson_plan": {
    "sections": [
      {"title": "教学目标", "paragraphs": ["..."]},
      {"title": "学情分析", "paragraphs": ["..."]},
      {"title": "教学重难点", "paragraphs": ["..."]},
      {"title": "教学过程与活动", "paragraphs": ["..."]},
      {"title": "作业与评价", "paragraphs": ["..."]}
    ]
  },
  "interactive": {
    "quiz": [
      {"question": "单选题题干", "options": ["A. ...", "B. ...", "C. ...", "D. ..."], "answer": "A"}
    ]
  }
}
要求：slides 至少 6 页、至多 15 页；每页 bullets 3～5 条；quiz 至少 3 题；全文中文。"""


def _rag_snippets_for_courseware(query: str, k: int = 5) -> str:
    if not rag_engine:
        return ""
    try:
        docs = rag_engine.retrieve(query, k=k)
        if not docs:
            return ""
        parts = []
        for i, d in enumerate(docs[:k]):
            meta = getattr(d, "metadata", None) or {}
            src = os.path.basename(str(meta.get("source", "material")))
            parts.append(f"[{i + 1} {src}] {d.page_content[:650]}")
        return "\n\n".join(parts)
    except Exception as ex:
        print(f"⚠️ RAG snippets for courseware: {ex}")
        return ""


def _build_courseware_payload_llm(
    final_requirements: str,
    user_outline: Optional[str],
    rag_context: str,
) -> Optional[Dict[str, Any]]:
    user_part = f"【教学需求】\n{final_requirements}\n"
    if user_outline and user_outline.strip():
        user_part += f"\n【用户给定大纲/要点】（请在此基础上扩展并结构化）\n{user_outline.strip()}\n"
    if rag_context:
        user_part += f"\n【参考资料片段】（可引用，勿编造事实）\n{rag_context}\n"
    data = llm_engine.generate_json(user_part, system_prompt=COURSEWARE_JSON_SYSTEM)
    if not data or data.get("fallback") is True or data.get("error"):
        return None
    slides = data.get("slides")
    if not isinstance(slides, list) or len(slides) < 1:
        return None
    return data


def _format_outline_from_payload(data: Dict[str, Any]) -> str:
    lines = [f"# {data.get('course_title', '课程')}"]
    for i, s in enumerate(data.get("slides", []), 1):
        if not isinstance(s, dict):
            continue
        lines.append(f"\n## 第{i}页 {s.get('title', '')}")
        for b in s.get("bullets") or []:
            lines.append(f"- {b}")
    return "\n".join(lines)


def _write_pptx_structured(path: Path, course_title: str, slides: List[Dict[str, Any]]):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide0 = prs.slides.add_slide(title_slide_layout)
    slide0.shapes.title.text = (course_title or "课程")[:120]
    slide0.placeholders[1].text = "AI 结构化生成（LLM）"

    bullet_layout = prs.slide_layouts[1]
    for s in slides[:25]:
        if not isinstance(s, dict):
            continue
        sld = prs.slides.add_slide(bullet_layout)
        sld.shapes.title.text = str(s.get("title") or "未命名")[:80]
        tf = sld.placeholders[1].text_frame
        tf.clear()
        bullets = s.get("bullets") or []
        if not bullets:
            bullets = ["（无要点）"]
        for i, point in enumerate(bullets[:8]):
            text = str(point)[:500]
            if i == 0:
                tf.text = text
            else:
                p = tf.add_paragraph()
                p.text = text
                p.level = 0
    prs.save(str(path))


def _write_docx_structured(path: Path, course_title: str, final_requirements: str, payload: Dict[str, Any]):
    doc = DocxDocument()
    doc.add_heading(f"{course_title} - 教案", level=1)
    doc.add_paragraph(f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    doc.add_heading("需求说明", level=2)
    doc.add_paragraph(final_requirements or "未提供")
    lp = payload.get("lesson_plan") or {}
    sections = lp.get("sections") if isinstance(lp, dict) else None
    if isinstance(sections, list) and sections:
        for sec in sections:
            if not isinstance(sec, dict):
                continue
            doc.add_heading(str(sec.get("title", "章节")), level=2)
            for para in sec.get("paragraphs") or []:
                doc.add_paragraph(str(para))
    else:
        doc.add_heading("课程大纲（来自幻灯片）", level=2)
        for i, s in enumerate(payload.get("slides", []), 1):
            if not isinstance(s, dict):
                continue
            doc.add_paragraph(f"{i}. {s.get('title', '')}", style="List Number")
            for b in s.get("bullets") or []:
                doc.add_paragraph(str(b), style="List Bullet")
    doc.save(str(path))


def _write_game_html_structured(path: Path, course_title: str, quiz: List[Dict[str, Any]]):
    items = []
    for idx, q in enumerate(quiz[:10], 1):
        if not isinstance(q, dict):
            continue
        opts = "".join(f"<li>{o}</li>" for o in (q.get("options") or []))
        ans = str(q.get("answer", "")).replace('"', "&quot;")
        items.append(
            f"""
<section class="q" data-ans="{ans}">
  <h3>第{idx}题</h3>
  <p>{q.get("question", "")}</p>
  <ul>{opts}</ul>
  <p><button type="button" onclick="check(this)">记录作答</button> <span class="tip"></span></p>
</section>"""
        )
    body = "\n".join(items) if items else "<p>暂无题目</p>"
    html = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{course_title} - 互动练习</title>
  <style>body{{font-family:sans-serif;max-width:640px;margin:1rem auto;}} .tip{{color:#080;margin-left:8px;}}</style>
</head>
<body>
  <h1>{course_title} - 互动练习</h1>
  <p>题目由 LLM 生成（演示页，可接判分与统计）。</p>
  {body}
  <script>
  function check(btn) {{
    var tip = btn.parentElement.querySelector('.tip');
    tip.textContent = '已记录';
  }}
  </script>
</body>
</html>"""
    path.write_text(html, encoding="utf-8")


def _build_outline(final_requirements: str, provided_outline: Optional[str]) -> str:
    if provided_outline and provided_outline.strip():
        return provided_outline.strip()
    prompt = (
        "请为以下教学需求生成一份结构化课程大纲（按章节列点，中文输出）：\n"
        f"{final_requirements or ''}"
    )
    system_prompt = "你是课程设计助手，请只输出清晰的大纲正文，不要输出额外解释。"
    return llm_engine.generate(prompt, system_prompt=system_prompt)


def _write_docx(path: Path, title: str, final_requirements: str, outline: str):
    doc = DocxDocument()
    doc.add_heading(f"{title} - 教案", level=1)
    doc.add_paragraph(f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    doc.add_heading("需求说明", level=2)
    doc.add_paragraph(final_requirements or "未提供")
    doc.add_heading("课程大纲", level=2)
    for line in outline.splitlines():
        text = line.strip()
        if text:
            doc.add_paragraph(text, style="List Bullet")
    doc.save(str(path))


def _write_pptx(path: Path, title: str, outline: str):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "大纲切分生成（降级模式）"

    bullet_layout = prs.slide_layouts[1]
    lines = [line.strip() for line in outline.splitlines() if line.strip()]
    chunks = [lines[i:i + 5] for i in range(0, len(lines), 5)] or [["暂无大纲内容"]]
    for idx, chunk in enumerate(chunks, start=1):
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = f"章节 {idx}"
        text_frame = s.placeholders[1].text_frame
        text_frame.clear()
        for i, point in enumerate(chunk):
            if i == 0:
                text_frame.text = point
            else:
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
    prs.save(str(path))


def _write_game_html(path: Path, title: str, outline: str):
    first_lines = [line.strip() for line in outline.splitlines() if line.strip()][:3]
    focus = "、".join(first_lines) if first_lines else "课程基础内容"
    html = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{title} - 互动练习</title>
</head>
<body>
  <h1>{title} - 互动练习</h1>
  <p>知识点：{focus}</p>
  <p>降级模式占位页；结构化生成成功时会替换为 AI 测验题。</p>
  <button onclick="alert('答对了！继续下一题。')">提交答案</button>
</body>
</html>"""
    path.write_text(html, encoding="utf-8")


@app.post("/api/courseware/generate", response_model=CoursewareGenerateResponse)
async def generate_courseware(request: CoursewareGenerateRequest):
    if not llm_engine:
        raise HTTPException(status_code=503, detail="LLM not initialized yet")
    if not request.sessionId.strip():
        raise HTTPException(status_code=400, detail="sessionId cannot be empty")
    if not request.finalRequirements.strip():
        raise HTTPException(status_code=400, detail="finalRequirements cannot be empty")

    task_id = f"py-task-{uuid4().hex[:12]}"
    folder = GENERATED_ROOT / _sanitize_filename(request.sessionId) / task_id
    folder.mkdir(parents=True, exist_ok=True)

    ppt_path = folder / "courseware.pptx"
    doc_path = folder / "teaching_plan.docx"
    game_path = folder / "interactive_game.html"

    llm_mock = bool(getattr(llm_engine, "is_mock_mode", True))
    rag_ctx = _rag_snippets_for_courseware(request.finalRequirements, k=5)
    payload = _build_courseware_payload_llm(request.finalRequirements, request.outline, rag_ctx)
    ai_structured = False

    if payload:
        ai_structured = True
        course_title = str(payload.get("course_title") or "AI课程内容")[:120]
        outline = _format_outline_from_payload(payload)
        slides = payload.get("slides") if isinstance(payload.get("slides"), list) else []
        _write_pptx_structured(ppt_path, course_title, slides)
        _write_docx_structured(doc_path, course_title, request.finalRequirements, payload)
        inter = payload.get("interactive") or {}
        quiz = inter.get("quiz") if isinstance(inter, dict) else None
        if not isinstance(quiz, list):
            quiz = []
        _write_game_html_structured(game_path, course_title, quiz)
    else:
        outline = _build_outline(request.finalRequirements, request.outline)
        title = "AI课程内容"
        _write_pptx(ppt_path, title, outline)
        _write_docx(doc_path, title, request.finalRequirements, outline)
        _write_game_html(game_path, title, outline)

    base = f"/generated-files/{_sanitize_filename(request.sessionId)}/{task_id}"
    return CoursewareGenerateResponse(
        taskId=task_id,
        status="success",
        progress=100,
        outline=outline,
        pptUrl=f"{base}/courseware.pptx",
        docUrl=f"{base}/teaching_plan.docx",
        gameUrl=f"{base}/interactive_game.html",
        model_used=SYSTEM_CONFIG.get("model_name", "unknown"),
        llm_mock=llm_mock,
        ai_structured=ai_structured,
    )


@app.get("/health")
async def health_check():
    mock = bool(getattr(llm_engine, "is_mock_mode", True)) if llm_engine else True
    return {
        "status": "healthy",
        "rag_ready": rag_engine is not None,
        "llm_ready": llm_engine is not None,
        "llm_mock": mock,
    }

if __name__ == "__main__":
    # 启动命令：uvicorn api_server:app --host 0.0.0.0 --port 8000 --reload
    uvicorn.run(app, host="0.0.0.0", port=8000)